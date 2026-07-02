from __future__ import annotations

from pathlib import Path

from docx import Document
from pypdf import PdfReader
from pptx import Presentation


SUPPORTED_EXTENSIONS = {".txt", ".md", ".pdf", ".docx", ".pptx"}
MAX_DOCUMENT_CHARS = 60_000


def extract_text(path: Path) -> str:
    ext = path.suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file type: {ext}")

    if ext in {".txt", ".md"}:
        text = path.read_text(encoding="utf-8", errors="replace")
    elif ext == ".pdf":
        text = _extract_pdf(path)
    elif ext == ".docx":
        text = _extract_docx(path)
    elif ext == ".pptx":
        text = _extract_pptx(path)
    else:
        text = ""

    text = "\n".join(line.rstrip() for line in text.splitlines())
    text = text.strip()
    if len(text) > MAX_DOCUMENT_CHARS:
        return text[:MAX_DOCUMENT_CHARS] + "\n\n[Document truncated because it is too long.]"
    return text


def _extract_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    parts = []
    for index, page in enumerate(reader.pages, start=1):
        page_text = page.extract_text() or ""
        if page_text.strip():
            parts.append(f"[Page {index}]\n{page_text}")
    return "\n\n".join(parts)


def _extract_docx(path: Path) -> str:
    doc = Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    table_parts = []
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                table_parts.append(" | ".join(cells))
    return "\n".join(paragraphs + table_parts)


def _extract_pptx(path: Path) -> str:
    presentation = Presentation(str(path))
    slides = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        shapes = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                shapes.append(shape.text.strip())
        if shapes:
            slides.append(f"[Slide {slide_index}]\n" + "\n".join(shapes))
    return "\n\n".join(slides)
